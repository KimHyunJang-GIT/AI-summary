"""extract module split from pipeline (auto-split from originals)."""
from __future__ import annotations
import re
import sys
from pathlib import Path
from typing import Tuple, Dict, Any

# -- Optional Dependencies --
try:
    import textract
except ImportError:
    textract = None

try:
    import win32com.client
except ImportError:
    win32com = None

try:
    import docx
except ImportError:
    docx = None

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    import pytesseract
except ImportError:
    pytesseract = None

try:
    import pptx
except ImportError:
    pptx = None


class TextCleaner:
    _multi_space = re.compile(r"\s+")
    @classmethod
    def clean(cls, s:str)->str:
        if not s: return ""
        s = s.replace("\x00"," ")
        s = cls._multi_space.sub(" ", s)
        return s.strip()

class BaseExtractor:
    exts: Tuple[str,...] = ()
    def can_handle(self, p:Path)->bool: return p.suffix.lower() in self.exts
    def extract(self, p:Path)->Dict[str,Any]: raise NotImplementedError

class HwpExtractor(BaseExtractor):
    exts=(".hwp",)
    def extract(self, p:Path)->Dict[str,Any]:
        if textract:
            try:
                t = textract.process(str(p)).decode("utf-8", "ignore")
                if t.strip():
                    return {"ok":True,"text":TextCleaner.clean(t)}
            except Exception as e:
                sys.stderr.write(f"[경고] textract HWP 추출 실패 ({p.name}): {e}\n")
        if win32com:
            sys.stderr.write(f"[정보] win32com HWP 추출은 복잡하여 현재 구현되어 있지 않습니다.\n")
        return {"ok":False,"text":"HWP 추출 실패"}

class DocxExtractor(BaseExtractor):
    exts=(".docx",".doc")
    def extract(self, p:Path)->Dict[str,Any]:
        if p.suffix.lower() == ".docx":
            if not docx: return {"ok": False, "text": "'python-docx' not installed"}
            try:
                d = docx.Document(str(p)); t = "\n".join(par.text for par in d.paragraphs)
                if t.strip(): return {"ok":True,"text":TextCleaner.clean(t)}
            except Exception as e: sys.stderr.write(f"[경고] python-docx DOCX 추출 실패 ({p.name}): {e}\n")
        if textract:
            try:
                t = textract.process(str(p)).decode("utf-8", "ignore")
                if t.strip(): return {"ok":True,"text":TextCleaner.clean(t)}
            except Exception as e: sys.stderr.write(f"[경고] textract DOC/DOCX 추출 실패 ({p.name}): {e}\n")
        if p.suffix.lower() == ".doc" and win32com:
            word_app = None
            try:
                word_app = win32com.client.Dispatch("Word.Application")
                word_app.Visible = False
                doc = word_app.Documents.Open(str(p), ConfirmConversions=False, ReadOnly=True)
                text = doc.Content.Text
                doc.Close(SaveChanges=False)
                if text.strip(): return {"ok":True,"text":TextCleaner.clean(text)}
            except Exception as e:
                sys.stderr.write(f"[경고] win32com DOC 추출 실패 ({p.name}): {e}\n")
            finally:
                if word_app:
                    word_app.Quit()
        return {"ok":False,"text":"DOC/DOCX 추출 실패"}

class PdfExtractor(BaseExtractor):
    exts=(".pdf",)
    def extract(self, p:Path)->Dict[str,Any]:
        if fitz:
            try:
                doc = fitz.open(p); text = "\n".join(page.get_text() for page in doc); doc.close()
                if text.strip(): return {"ok":True,"text":TextCleaner.clean(text)}
            except Exception as e: sys.stderr.write(f"[경고] PyMuPDF PDF 추출 실패 ({p.name}): {e}\n")
        if pytesseract:
            try:
                pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
                doc = fitz.open(p)
                text_parts = []
                for page_num in range(doc.page_count):
                    page = doc[page_num]
                    pix = page.get_pixmap()
                    img_bytes = pix.tobytes("png")
                    text_parts.append(pytesseract.image_to_string(img_bytes, lang='kor+eng'))
                doc.close()
                text = "\n".join(text_parts)
                if text.strip(): return {"ok":True,"text":TextCleaner.clean(text)}
            except Exception as e: sys.stderr.write(f"[경고] Tesseract OCR PDF 추출 실패 ({p.name}): {e}\n")
        if textract:
            try:
                t = textract.process(str(p)).decode("utf-8", "ignore")
                if t.strip(): return {"ok":True,"text":TextCleaner.clean(t)}
            except Exception as e: sys.stderr.write(f"[경고] textract PDF 추출 실패 ({p.name}): {e}\n")
        return {"ok":False,"text":"PDF 추출 실패"}

class PptxExtractor(BaseExtractor):
    exts=(".pptx",".ppt")
    def extract(self, p:Path)->Dict[str,Any]:
        if p.suffix.lower() == ".pptx":
            if not pptx: return {"ok": False, "text": "'python-pptx' not installed"}
            try:
                pres = pptx.Presentation(str(p)); texts = [shape.text for slide in pres.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                if texts: return {"ok":True,"text":TextCleaner.clean("\n".join(texts))}
            except Exception as e: sys.stderr.write(f"[경고] python-pptx PPTX 추출 실패 ({p.name}): {e}\n")
        if textract:
            try:
                t = textract.process(str(p)).decode("utf-8", "ignore")
                if t.strip(): return {"ok":True,"text":TextCleaner.clean(t)}
            except Exception as e: sys.stderr.write(f"[경고] textract PPT/PPTX 추출 실패 ({p.name}): {e}\n")
        if p.suffix.lower() == ".ppt" and win32com:
            powerpoint_app = None
            try:
                powerpoint_app = win32com.client.Dispatch("Powerpoint.Application")
                powerpoint_app.Visible = False
                pres = powerpoint_app.Presentations.Open(str(p), ReadOnly=True, WithWindow=False)
                text_parts = []
                for slide in pres.Slides:
                    for shape in slide.Shapes:
                        if shape.HasTextFrame:
                            text_parts.append(shape.TextFrame.TextRange.Text)
                pres.Close()
                text = "\n".join(text_parts)
                if text.strip(): return {"ok":True,"text":TextCleaner.clean(text)}
            except Exception as e:
                sys.stderr.write(f"[경고] win32com PPT 추출 실패 ({p.name}): {e}\n")
            finally:
                if powerpoint_app:
                    powerpoint_app.Quit()
        return {"ok":False,"text":"PPT/PPTX 추출 실패"}


EXT_MAP = {
    ext: extractor_class()
    for extractor_class in [HwpExtractor, DocxExtractor, PdfExtractor, PptxExtractor]
    for ext in extractor_class.exts
}
