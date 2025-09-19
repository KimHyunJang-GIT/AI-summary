# -*- coding: utf-8 -*-
"""
Step2: 텍스트 추출(CorpusBuilder) + 인덱싱 실행(run_indexing)
- HWP 추출 기능 추가 (textract)
- PDF 추출 성능 향상 (PyMuPDF + Tesseract OCR + textract)
- Excel/CSV 추출 성능 향상 (인코덱, 라이브러리 체크)
- DOC/PPT 구형 형식 추출 기능 추가 (textract, win32com)
- 모든 문서에 대한 자동 요약 기능 추가
"""
import os, re, sys, time, threading
from pathlib import Path
from dataclasses import dataclass
from typing import Optional, Dict, Any, List, Tuple

# --- 의존성 라이브러리 ---
try:
    import pandas as pd
except ImportError: pd = None
try:
    import fitz  # PyMuPDF
except ImportError: fitz = None
try:
    import docx
except ImportError: docx = None
try:
    import pptx
except ImportError: pptx = None
try:
    import pytesseract # Tesseract OCR 파이썬 래퍼
except ImportError: pytesseract = None
try:
    import textract # 다양한 문서 형식 지원
except ImportError: textract = None
try:
    import win32com.client # Windows COM (MS Office 연동)
except ImportError: win32com = None

try:
    from tqdm import tqdm
except ImportError: tqdm = None
try:
    from sentence_transformers import SentenceTransformer, util
except ImportError: SentenceTransformer, util = None, None

from retriever import Retriever

__all__ = ["CorpusBuilder", "run_indexing"]

# =========================
# 유틸리티 함수
# =========================
class TextCleaner:
    _multi_space = re.compile(r"\s+")
    @classmethod
    def clean(cls, s:str)->str:
        if not s: return ""
        s = s.replace("\x00"," ")
        s = cls._multi_space.sub(" ", s)
        return s.strip()

def simple_summary(text: str, max_len: int = 400) -> str:
    if not text: return ""
    lines = [line.strip() for line in text.split('\n') if line.strip()]

    if lines:
        summary_parts = lines[:3]
        summary = " ".join(summary_parts)
    else:
        summary = text

    return (summary[:max_len] + "…") if len(summary) > max_len else summary

# =========================
# 파일 유형별 텍스트 추출기
# =========================
class BaseExtractor:
    exts: Tuple[str,...] = ()
    def can_handle(self, p:Path)->bool: return p.suffix.lower() in self.exts
    def extract(self, p:Path)->Dict[str,Any]: raise NotImplementedError

class HwpExtractor(BaseExtractor):
    exts=(".hwp",)
    def extract(self, p:Path)->Dict[str,Any]:
        # 1. textract로 HWP 추출 시도
        if textract:
            try:
                t = textract.process(str(p)).decode("utf-8", "ignore")
                if t.strip():
                    return {"ok":True,"text":TextCleaner.clean(t)}
            except Exception as e:
                sys.stderr.write(f"[경고] textract HWP 추출 실패 ({p.name}): {e}\n")

        # 2. win32com을 통한 HWP 추출 (한글 오피스 설치 필요, 매우 환경 의존적)
        #    HWP COM 객체 연동은 복잡하며, 설치된 한글 오피스 버전에 따라 API가 다를 수 있습니다.
        #    textract가 작동하지 않는다면, 'hwp5.py' 또는 'pyhwp'와 같은 외부 라이브러리 사용을 고려해야 합니다.
        if win32com:
            sys.stderr.write(f"[정보] win32com HWP 추출은 복잡하여 현재 구현되어 있지 않습니다. textract 실패 시 다른 HWP 라이브러리를 고려하세요.\n")

        return {"ok":False,"text":"HWP 추출 실패"}

class DocxExtractor(BaseExtractor):
    exts=(".docx",".doc") # .doc 추가
    def extract(self, p:Path)->Dict[str,Any]:
        # 1. python-docx로 .docx 추출 시도
        if p.suffix.lower() == ".docx":
            if not docx: return {"ok": False, "text": "'python-docx' not installed"}
            try:
                d = docx.Document(str(p)); t = "\n".join(par.text for par in d.paragraphs)
                if t.strip(): return {"ok":True,"text":TextCleaner.clean(t)}
            except Exception as e: sys.stderr.write(f"[경고] python-docx DOCX 추출 실패 ({p.name}): {e}\n")

        # 2. textract로 .doc 또는 .docx 추출 시도
        if textract:
            try:
                t = textract.process(str(p)).decode("utf-8", "ignore")
                if t.strip(): return {"ok":True,"text":TextCleaner.clean(t)}
            except Exception as e: sys.stderr.write(f"[경고] textract DOC/DOCX 추출 실패 ({p.name}): {e}\n")

        # 3. win32com으로 .doc 추출 시도 (MS Word 설치 필요)
        if p.suffix.lower() == ".doc" and win32com:
            word_app = None
            try:
                word_app = win32com.client.Dispatch("Word.Application")
                word_app.Visible = False # Word 애플리케이션을 숨김
                doc = word_app.Documents.Open(str(p), ConfirmConversions=False, ReadOnly=True)
                text = doc.Content.Text
                doc.Close(SaveChanges=False) # 문서 닫기 (변경사항 저장 안 함)
                if text.strip(): return {"ok":True,"text":TextCleaner.clean(text)}
            except Exception as e:
                sys.stderr.write(f"[경고] win32com DOC 추출 실패 ({p.name}): {e}\n")
            finally:
                if word_app:
                    word_app.Quit() # Word 애플리케이션 종료

        return {"ok":False,"text":"DOC/DOCX 추출 실패"}

class ExcelExtractor(BaseExtractor):
    exts=(".xlsx",".xls",".xlsm",".xlsb",".xltx",".csv")
    def extract(self, p:Path)->Dict[str,Any]:
        if not pd: return {"ok": False, "text": "'pandas' not installed"}

        # CSV 파일 처리
        if p.suffix.lower() == ".csv":
            for enc in ['utf-8', 'cp949', 'euc-kr']:
                try:
                    df = pd.read_csv(p, nrows=200, on_bad_lines='skip', encoding=enc)
                    header = " | ".join(map(str, df.columns)); rows_text = "\n".join([" | ".join(map(str, row)) for _, row in df.head(20).iterrows()])
                    return {"ok":True,"text":TextCleaner.clean(f"{header}\n{rows_text}"), "columns": df.columns.tolist()}
                except Exception as e:
                    sys.stderr.write(f"[경고] CSV 추출 실패 ({p.name}) with {enc}: {e}\n")
            return {"ok":False,"text": f"CSV 추출 실패 (인코딩 문제)"}

        # Excel 파일 처리
        else: # .xlsx, .xls, .xlsm, .xlsb, .xltx
            # openpyxl (xlsx, xlsm, xltx)
            if p.suffix.lower() in (".xlsx", ".xlsm", ".xltx"):
                try:
                    df = pd.read_excel(p, sheet_name=0, nrows=200, engine='openpyxl')
                    header = " | ".join(map(str, df.columns)); rows_text = "\n".join([" | ".join(map(str, row)) for _, row in df.head(20).iterrows()])
                    return {"ok":True, "text":TextCleaner.clean(f"{header}\n{rows_text}"), "columns": df.columns.tolist()}
                except ImportError: sys.stderr.write("[경고] 'openpyxl' not installed. Excel .xlsx/.xlsm/.xltx files may fail.\n")
                except Exception as e: sys.stderr.write(f"[경고] openpyxl Excel 추출 실패 ({p.name}): {e}\n")

            # pyxlsb (xlsb)
            elif p.suffix.lower() == ".xlsb":
                try:
                    df = pd.read_excel(p, sheet_name=0, nrows=200, engine='pyxlsb')
                    header = " | ".join(map(str, df.columns)); rows_text = "\n".join([" | ".join(map(str, row)) for _, row in df.head(20).iterrows()])
                    return {"ok":True, "text":TextCleaner.clean(f"{header}\n{rows_text}"), "columns": df.columns.tolist()}
                except ImportError: sys.stderr.write("[경고] 'pyxlsb' not installed. Excel .xlsb files may fail.\n")
                except Exception as e: sys.stderr.write(f"[경고] pyxlsb Excel 추출 실패 ({p.name}): {e}\n")

            # xlrd (xls)
            elif p.suffix.lower() == ".xls":
                try:
                    df = pd.read_excel(p, sheet_name=0, nrows=200, engine='xlrd')
                    header = " | ".join(map(str, df.columns)); rows_text = "\n".join([" | ".join(map(str, row)) for _, row in df.head(20).iterrows()])
                    return {"ok":True, "text":TextCleaner.clean(f"{header}\n{rows_text}"), "columns": df.columns.tolist()}
                except ImportError: sys.stderr.write("[경고] 'xlrd' not installed. Excel .xls files may fail.\n")
                except Exception as e: sys.stderr.write(f"[경고] xlrd Excel 추출 실패 ({p.name}): {e}\n")

            # textract (Excel 폴백)
            if textract:
                try:
                    t = textract.process(str(p)).decode("utf-8", "ignore")
                    if t.strip(): return {"ok":True,"text":TextCleaner.clean(t)}
                except Exception as e: sys.stderr.write(f"[경고] textract Excel 추출 실패 ({p.name}): {e}\n")

            # win32com (Excel 폴백 - MS Excel 설치 필요) - 구현 필요 시 추가
            if win32com:
                try:
                    # Excel COM 객체 생성 및 텍스트 추출 (구현 필요)
                    pass # 현재는 구현되어 있지 않음
                except Exception as e: sys.stderr.write(f"[경고] win32com Excel 추출 실패 ({p.name}): {e}\n")

            return {"ok":False,"text": f"Excel 추출 실패 (지원되지 않는 형식 또는 오류)"}

class PdfExtractor(BaseExtractor):
    exts=(".pdf",)
    def extract(self, p:Path)->Dict[str,Any]:
        # 1. PyMuPDF로 텍스트 추출 시도
        if fitz:
            try:
                doc = fitz.open(p); text = "\n".join(page.get_text() for page in doc); doc.close()
                if text.strip(): return {"ok":True,"text":TextCleaner.clean(text)}
            except Exception as e: sys.stderr.write(f"[경고] PyMuPDF PDF 추출 실패 ({p.name}): {e}\n")

        # 2. PyMuPDF 실패 시, Tesseract OCR 시도
        if pytesseract:
            try:
                # Tesseract OCR 경로 설정 (설치 경로에 맞게 수정 필요)
                pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

                doc = fitz.open(p)
                text_parts = []
                for page_num in range(doc.page_count):
                    page = doc[page_num]
                    pix = page.get_pixmap()
                    img_bytes = pix.tobytes("png")
                    text_parts.append(pytesseract.image_to_string(img_bytes, lang='kor+eng'))
                doc.close()
                text = "\n".join(text_parts)
                if text.strip(): return {"ok":True,"text":TextCleaner.clean(text)}
            except Exception as e: sys.stderr.write(f"[경고] Tesseract OCR PDF 추출 실패 ({p.name}): {e}\n")

        # 3. 모든 시도 실패 시, textract 시도
        if textract:
            try:
                t = textract.process(str(p)).decode("utf-8", "ignore")
                if t.strip(): return {"ok":True,"text":TextCleaner.clean(t)}
            except Exception as e: sys.stderr.write(f"[경고] textract PDF 추출 실패 ({p.name}): {e}\n")

        return {"ok":False,"text":"PDF 추출 실패"}

class PptxExtractor(BaseExtractor):
    exts=(".pptx",".ppt") # .ppt 추가
    def extract(self, p:Path)->Dict[str,Any]:
        # 1. python-pptx로 .pptx 추출 시도
        if p.suffix.lower() == ".pptx":
            if not pptx: return {"ok": False, "text": "'python-pptx' not installed"}
            try:
                pres = pptx.Presentation(str(p)); texts = [shape.text for slide in pres.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                if texts: return {"ok":True,"text":TextCleaner.clean("\n".join(texts))}
            except Exception as e: sys.stderr.write(f"[경고] python-pptx PPTX 추출 실패 ({p.name}): {e}\n")

        # 2. textract로 .ppt 또는 .pptx 추출 시도
        if textract:
            try:
                t = textract.process(str(p)).decode("utf-8", "ignore")
                if t.strip(): return {"ok":True,"text":TextCleaner.clean(t)}
            except Exception as e: sys.stderr.write(f"[경고] textract PPT/PPTX 추출 실패 ({p.name}): {e}\n")

        # 3. win32com으로 .ppt 추출 시도 (MS PowerPoint 설치 필요)
        if p.suffix.lower() == ".ppt" and win32com:
            powerpoint_app = None
            try:
                powerpoint_app = win32com.client.Dispatch("Powerpoint.Application")
                powerpoint_app.Visible = False # PowerPoint 애플리케이션을 숨김
                pres = powerpoint_app.Presentations.Open(str(p), ReadOnly=True, WithWindow=False)
                text_parts = []
                for slide in pres.Slides:
                    for shape in slide.Shapes:
                        if shape.HasTextFrame:
                            text_parts.append(shape.TextFrame.TextRange.Text)
                pres.Close() # 프레젠테이션 닫기
                text = "\n".join(text_parts)
                if text.strip(): return {"ok":True,"text":TextCleaner.clean(text)}
            except Exception as e:
                sys.stderr.write(f"[경고] win32com PPT 추출 실패 ({p.name}): {e}\n")
            finally:
                if powerpoint_app:
                    powerpoint_app.Quit() # PowerPoint 애플리케이션 종료

        return {"ok":False,"text":"PPT/PPTX 추출 실패"}

EXTRACTORS = [HwpExtractor(), DocxExtractor(), ExcelExtractor(), PdfExtractor(), PptxExtractor()]
EXT_MAP = {e:ex for ex in EXTRACTORS for e in ex.exts}

# =========================
# 코퍼스 빌더 (AI 요약 기능 추가)
# =========================
@dataclass
class ExtractRecord:
    path:str; ext:str; ok:bool; text:str; summary:str; title:str; columns:List[str]

class CorpusBuilder:
    TOPIC_CANDIDATES = [
        "파이썬 프로그래밍", "자바 프로그래밍", "SQL 데이터베이스", "웹 개발", "자바스크립트",
        "머신러닝 및 인공지능", "데이터 분석 및 시각화", "클라우드 컴퓨팅", "보안",
        "업무 보고서", "회의록", "기획서", "제안서", "계약서", "법률 문서", "인사 관리", "재무 및 회계",
        "강의 교안", "연구 논문", "기술 문서", "사용자 매뉴얼",
    ]

    def __init__(self, max_text_chars:int=500_000, progress:bool=True):
        self.max_text_chars=max_text_chars
        self.progress=progress
        if SentenceTransformer is None:
            raise ImportError("sentence-transformers is not installed. Please run: pip install sentence-transformers")
        print("🧠 Loading model for summarization...")
        self.semantic_model = SentenceTransformer(Retriever.MODEL_NAME)
        self.topic_embeddings = self.semantic_model.encode(self.TOPIC_CANDIDATES, convert_to_tensor=True)

    def _generate_ai_summary(self, text: str) -> str:
        if not text or not text.strip():
            return "(내용이 없어 요약할 수 없습니다)"
        try:
            doc_embedding = self.semantic_model.encode(text, convert_to_tensor=True)
            cos_scores = util.cos_sim(doc_embedding, self.topic_embeddings)[0]
            best_topic_index = cos_scores.argmax()
            best_topic = self.TOPIC_CANDIDATES[best_topic_index]
            return f"이 문서는 '{best_topic}' 관련 자료로 보입니다."
        except Exception as e:
            return f"[요약 오류: {e}]"

    def build(self, file_rows:List[Dict[str,Any]]) -> pd.DataFrame:
        if not pd: raise RuntimeError("pandas is required. Please run: pip install pandas")
        iterator = tqdm(file_rows, desc="📥 Extracting & Summarizing", unit="file") if self.progress and tqdm else file_rows
        records: List[ExtractRecord] = []

        for row in iterator:
            p = Path(row["path"]); ext = p.suffix.lower()
            extractor = EXT_MAP.get(ext)
            if not extractor:
                records.append(ExtractRecord(str(p), ext, False, "", "(알 수 없는 파일 형식)", p.stem, [])); continue

            try:
                extract_result = extractor.extract(p)
                raw_text = (extract_result.get("text","") or "")[:self.max_text_chars]

                # --- 디버깅을 위한 raw_text 출력 ---
                if not raw_text.strip():
                    sys.stderr.write(f"[디버그] '{p.name}'에서 추출된 텍스트가 비어있습니다. 원본 파일 확인 필요.\n")
                elif len(raw_text) < 50: # 너무 짧은 텍스트도 의심
                    sys.stderr.write(f"[디버그] '{p.name}'에서 추출된 텍스트가 너무 짧습니다 ({len(raw_text)}자): '{raw_text}'\n")
                # ----------------------------------

                summary = self._generate_ai_summary(raw_text)
                columns = extract_result.get("columns", [])
                records.append(ExtractRecord(str(p), ext, bool(extract_result.get("ok",False)), raw_text, summary, p.stem, columns))
            except Exception as e:
                records.append(ExtractRecord(str(p), ext, False, f"Build Error: {e}", "(처리 중 오류 발생)", p.stem, []))

        df = pd.DataFrame(records)
        print(f"✅ Text extraction & summarization complete: {int(df['ok'].sum())} successful, {len(df)} total.")
        return df

    @staticmethod
    def save(df: pd.DataFrame, out_path: Path):
        out_path.parent.mkdir(parents=True, exist_ok=True)
        try:
            df.to_parquet(out_path, index=False); print(f"💾 Corpus saved to Parquet: {out_path}")
        except Exception as e:
            csv_path = out_path.with_suffix(".csv"); df.to_csv(csv_path, index=False, encoding="utf-8")
            print(f"⚠️ Parquet engine not found, saved to CSV instead: {csv_path}\n   Error: {e}")

# =========================
# 인덱싱 실행 함수
# =========================
def run_indexing(corpus_path: Path, cache_dir: Path):
    print("🚀 Starting semantic indexing...")
    retriever = Retriever(corpus_path=corpus_path, cache_dir=cache_dir)
    retriever.ready(rebuild=True)
    print("✨ Indexing complete.")
